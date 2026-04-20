"""Format-specific text extractors (lazy-import optional dependencies)."""

from __future__ import annotations

import csv
import json
import logging
import xml.etree.ElementTree as ET
from collections.abc import Callable
from pathlib import Path

logger = logging.getLogger(__name__)

# Import name from ModuleNotFoundError → ai-builder optional extra (see pyproject.toml).
_OPTIONAL_IMPORT_TO_EXTRA: dict[str, str] = {
    "pdfplumber": "docs-pdf",
    "docx": "docs-word",
    "pptx": "docs-slides",
    "bs4": "docs-html",
    "striprtf": "docs-rtf",
    "openpyxl": "docs-spreadsheet",
    "xlrd": "docs-spreadsheet",
    "ebooklib": "docs-epub",
    "odf": "docs-odt",
}


def _missing_optional_dep_message(exc: ModuleNotFoundError) -> str:
    name = exc.name or ""
    extra = _OPTIONAL_IMPORT_TO_EXTRA.get(name, "docs")
    return (
        f"Missing optional dependency {name!r}. "
        f'Install with: pip install "ai-builder[{extra}]" '
        f'(or pip install "ai-builder[docs]" for all document formats).'
    )


# Extensions the umbrella loader accepts (single-family tools use subsets).
DEFAULT_UMBRELLA_FORMATS: list[str] = [
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".pdf",
    ".docx",
    ".doc",
    ".dot",
    ".dotx",
    ".pptx",
    ".ppt",
    ".html",
    ".htm",
    ".rtf",
    ".xlsx",
    ".xls",
    ".epub",
    ".odt",
]


def extract_plain(path: Path) -> str:
    """Plain text / markdown."""
    return path.read_text(encoding="utf-8", errors="replace")


def extract_csv(path: Path) -> str:
    lines: list[str] = []
    with path.open(newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f):
            lines.append("\t".join(row))
    return "\n".join(lines)


def extract_json(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    data = json.loads(raw)
    return json.dumps(data, ensure_ascii=False, indent=2)


def extract_xml(path: Path) -> str:
    tree = ET.parse(path)
    parts = [t.strip() for t in tree.getroot().itertext() if t and t.strip()]
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        return "\n\n".join(p.extract_text() or "" for p in pdf.pages)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_doc_legacy(path: Path) -> str:
    import docx2txt

    return docx2txt.process(str(path))


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n\n".join(texts)


def extract_ppt_unavailable(_path: Path) -> str:
    raise NotImplementedError(
        "Legacy .ppt (PowerPoint 97–2003) is not extracted in-process. "
        "Convert to .pptx or use an external converter."
    )


def extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    html = path.read_text(encoding="utf-8", errors="replace")
    return BeautifulSoup(html, "html.parser").get_text(separator="\n")


def extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                lines.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(lines)


def extract_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(str(path))
    lines: list[str] = []
    for sheet in wb.sheets():
        for row_idx in range(sheet.nrows):
            row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
            lines.append("\t".join(row))
    return "\n".join(lines)


def extract_epub(path: Path) -> str:
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(str(path))
    chunks: list[str] = []
    for item in book.get_items():
        if item.get_type() != ebooklib.ITEM_DOCUMENT:
            continue
        body = item.get_content()
        if not body:
            continue
        chunks.append(
            BeautifulSoup(body, "html.parser").get_text(separator="\n"),
        )
    return "\n\n".join(chunks)


def extract_odt(path: Path) -> str:
    from odf.opendocument import load
    from odf import teletype
    from odf.text import P

    doc = load(str(path))
    return "\n".join(teletype.extractText(p) for p in doc.getElementsByType(P))


_SUFFIX_HANDLERS: dict[str, Callable[[Path], str]] = {
    ".txt": extract_plain,
    ".md": extract_plain,
    ".csv": extract_csv,
    ".json": extract_json,
    ".xml": extract_xml,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".dotx": extract_docx,
    ".doc": extract_doc_legacy,
    ".dot": extract_doc_legacy,
    ".pptx": extract_pptx,
    ".ppt": extract_ppt_unavailable,
    ".html": extract_html,
    ".htm": extract_html,
    ".rtf": extract_rtf,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".epub": extract_epub,
    ".odt": extract_odt,
}


def extract_for_path(path: Path, allowed_suffixes: set[str] | None = None) -> str:
    """Extract text from path; optional filter by lowercase extension (with dot)."""
    suf = path.suffix.lower()
    if allowed_suffixes is not None and suf not in allowed_suffixes:
        return ""
    handler = _SUFFIX_HANDLERS.get(suf)
    if handler is None:
        return ""
    try:
        return handler(path)
    except NotImplementedError:
        raise
    except ModuleNotFoundError as exc:
        raise RuntimeError(_missing_optional_dep_message(exc)) from exc
    except Exception:
        logger.warning("Extractor failed for %s", path, exc_info=True)
        raise


def suffixes_supported() -> frozenset[str]:
    return frozenset(_SUFFIX_HANDLERS.keys())
