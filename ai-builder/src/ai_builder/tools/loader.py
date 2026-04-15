"""Document loader — reads files from a directory with format-specific extractors.

Supported formats (all extractors use lazy imports so only stdlib is required at
import time — heavy libraries are loaded on first use of that format):

  No extra deps:  .txt, .md, .csv, .json, .xml
  pdfplumber:     .pdf
  python-docx:    .docx
  docx2txt:       .doc   (legacy Word)
  python-pptx:    .pptx
  beautifulsoup4: .html, .htm
  striprtf:       .rtf
  openpyxl:       .xlsx
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

from ai_builder.core.tool import BaseTool, ToolInput, ToolOutput

logger = logging.getLogger(__name__)


class LoaderConfig(BaseModel):
    """Configuration for the document loader."""

    source_dir: str = Field(default="data/raw", description="Directory to load documents from")
    supported_formats: list[str] = Field(
        default=[
            ".txt", ".md", ".csv", ".json", ".xml",
            ".pdf", ".docx", ".doc", ".pptx",
            ".html", ".htm", ".rtf", ".xlsx",
        ],
        description="File extensions to process",
    )
    recursive: bool = Field(default=True, description="Recurse into subdirectories")

    model_config = {"extra": "allow"}


class LoaderInput(ToolInput):
    """Input: path to a directory or single file."""

    data: str = Field(default="data/raw", description="Source directory or file path")


class LoaderOutput(ToolOutput):
    """Output: list of {text, source, filename, format} dicts."""

    data: list[dict[str, Any]] = Field(default_factory=list)


class DocumentLoader(BaseTool[LoaderInput, LoaderOutput]):
    """
    Load and extract text from documents.

    Supports: TXT, MD, CSV, JSON, XML, PDF, DOCX, DOC, PPTX, HTML, HTM,
    RTF, XLSX.  Heavy libraries (pdfplumber, openpyxl, …) are imported
    lazily — only when a file of that type is actually processed.
    """

    name = "loader"
    description = "Load documents from filesystem with format-specific text extraction"
    version = "1.1.0"

    def __init__(self, config: LoaderConfig | None = None) -> None:
        self.config = config or LoaderConfig()

    def execute(self, inp: LoaderInput) -> LoaderOutput:
        source = Path(inp.data) if inp.data else Path(self.config.source_dir)
        if not source.exists():
            return LoaderOutput(data=[], success=False, error=f"Path not found: {source}")

        files = self._collect_files(source)
        docs: list[dict[str, Any]] = []
        errors: list[str] = []

        for f in files:
            try:
                text = self._extract(f)
                if text:
                    docs.append({
                        "text": text,
                        "source": str(f),
                        "filename": f.name,
                        "format": f.suffix.lower().lstrip("."),
                        "chars": len(text),
                    })
            except Exception as exc:
                errors.append(f"{f.name}: {exc}")
                logger.warning("Failed to extract %s: %s", f, exc)

        return LoaderOutput(
            data=docs,
            metadata={
                **inp.metadata,
                "file_count": len(docs),
                "error_count": len(errors),
                "source_dir": str(source),
                "errors": errors,
            },
        )

    def _collect_files(self, source: Path) -> list[Path]:
        suffixes = set(self.config.supported_formats)
        if source.is_file():
            return [source] if source.suffix.lower() in suffixes else []
        glob_fn = source.rglob if self.config.recursive else source.glob
        return sorted(f for f in glob_fn("*") if f.is_file() and f.suffix.lower() in suffixes)

    @staticmethod
    def _extract(path: Path) -> str:
        suffix = path.suffix.lower()

        # ── Plain text (stdlib only) ──
        if suffix in (".txt", ".csv", ".json", ".xml", ".md"):
            return path.read_text(encoding="utf-8", errors="replace")

        # ── PDF ──
        if suffix == ".pdf":
            import pdfplumber

            with pdfplumber.open(path) as pdf:
                return "\n\n".join(p.extract_text() or "" for p in pdf.pages)

        # ── Word DOCX (Open XML) ──
        if suffix == ".docx":
            from docx import Document

            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)

        # ── Word DOC (legacy binary) ──
        if suffix == ".doc":
            import docx2txt

            return docx2txt.process(str(path))

        # ── PowerPoint PPTX ──
        if suffix == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n\n".join(texts)

        # ── HTML / HTM ──
        if suffix in (".html", ".htm"):
            from bs4 import BeautifulSoup

            html = path.read_text(encoding="utf-8", errors="replace")
            return BeautifulSoup(html, "html.parser").get_text(separator="\n")

        # ── RTF ──
        if suffix == ".rtf":
            from striprtf.striprtf import rtf_to_text

            return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))

        # ── Excel XLSX ──
        if suffix == ".xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(str(path), read_only=True, data_only=True)
            lines: list[str] = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    lines.append("\t".join(cells))
            wb.close()
            return "\n".join(lines)

        return ""
